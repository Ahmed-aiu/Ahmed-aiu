import os, re
from docx import Document
from pptx import Presentation
from openai import AzureOpenAI
import pymupdf



def extract_text_from_docx(file):
    # Extract text from all docx pages and return as a single string.
    doc = Document(file)
    return ' '.join([para.text for para in doc.paragraphs])

def extract_text_from_pdf(file_obj):
    """
    Extracts text from a PDF file using PyMuPDF.
    :param file_obj: A file object or file-like object.
    :return: Extracted text from the PDF.
    """
    try:
        # Open the PDF from the file object (use "rb" mode for binary reading)
        doc = pymupdf.open(stream=file_obj.read(), filetype="pdf")
        
        text = ""
        # Loop through each page and extract text
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)  # Get the page
            text += page.get_text()  # Extract the text

        return text
    except Exception as e:
        return f"Error while extracting text: {str(e)}"

def extract_text_from_pptx(file):
    # Extract text from all pptx slides and return as a single string.
    presentation = Presentation(file)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_file(file):
    # Determine the file type and call the appropriate function to extract text.
    _,file_extension = os.path.splitext(file.filename)
    if file_extension == '.docx':
        return extract_text_from_docx(file)
    elif file_extension == '.pdf':
        return extract_text_from_pdf(file)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file)
    else:
        return "Unsupported file type"
    
def extract_text_from_files(files):
    text_files = ""
    for filename, file_storage in files.items():
        text_file = extract_text_from_file(file_storage)
        text_files += text_file + "\n"  # Adding a newline for separation between texts
    return text_files

def analysis(files): 
    # Access the environment variables from .env file
    azure_oai_endpoint = os.getenv('azure_oai_endpoint')
    azure_oai_key = os.getenv('azure_oai_key')
    azure_oai_deployment = os.getenv('azure_oai_deployment')
    # Extract text from the file
    textRFQ = extract_text_from_files(files)

    # Interact with extracted text with Azure OpenAI client
    try: 
        # Initialize the Azure OpenAI client
        client = AzureOpenAI(
         azure_endpoint = azure_oai_endpoint, 
         api_key=azure_oai_key,  
         api_version="2024-02-15-preview"
         )

        # # Define Instructions (old)
        # instructions = """You are a Sales-Assistant of an automotive engineering service provider company trained to assign business requests (RFQs) to responsible technical departments.
        #                 To do so you have to:
        #                 1. Learn the department structure (DEPNAMES), defining what department is responsible for what categories
        #                 2. Learn what department is responbile for what categories (DEPSTRUCT)
        #                 3. Summarize the RFQ text into a SUMMARY
        #                 4. Identify the categories (SYSTEMS) in the RFQ text 
        #                 5. Identify the work activities (DEVDIS) in the RFQ text 
        #                 5. Identify the responsible departments 
        #                 6. Evaluate the share of responsibility of each department based on the occurance of the categories, the departments are responsible for
        #                 7. Identify the Project start if not identifiable please state that (PSTART)
        #                 8. Identify the Project duration if not identifiable please state that (PDURATION)
        #                 9. Identify the date until when the proposal shall be submitted (DEADLINE)
        #                 10. Identify the contact email and or person for technical questions if not identifiable please state that (CONTACTTECH)
        #                 11. Identify the contact email and or person for commercial questions (maybe not a Person but a generic contact mail) if not identifiable please state 

        #                 ### 1. Learn the department structure (DEPLIST) ###
        #                 A department name is made up of three letters, followed by a dash, follwed by another letter. example: "BES-C". 
        #                 "BES" is the "System Department"
        #                 "BEB" is the "Battery Department"
        #                 "BEM" is the "E-Motor and powerelectronics Department"

        #                 "-C": "Function Developmnt and Controls"
        #                 "-C": "Systems Engineering" 
        #                 "-M": "Mechanical Engineering" 
        #                 "-V": "Validation and Verification" 
        #                 "-V": "Calibration" 
        #                 "-T": "Testing"

        #                 So "BES-C" is the "System Department, Function Development and Controls" and "BEM-V": is the "E-Motor and Powerelectronics Department, Calibration"
        #                 ###

        #                 ### 2. Learn what department (DEPSTRUCT) is responbile for what categories  ###
        #                 Learn the following responsible departments for the categories (SYSTEMS) and the work activivities (DEVDIS) and the name of the contactperson (CONTACT). the JSON-syntax is:
        #                 "DEPARTMENT":{"SYSTEMS":["System 1","System 1"],"DEVDIS":["work activity 1","work activity 2"],"CONTACT":"contact-person"}
        #                 example: "BES":{"SYSTEMS":["Powertrain and Charging System","Energy System"],"DEVDIS":["Project Management","Supplier Handling"],"CONTACT":"Dr. Felix Richert"}
        #                 means: "The Department BES is respondible for the categories/SYSTEMS Powertrain and Charging System and Energy System and is responible for the work activities (DEVDIS) Project Management and Supplier Handling and Dr. Felix Richert is the contact-person."
        #                 Here the (DEPSTRUCT):
        #                 {"BES":{"SYSTEMS":["Powertrain and Charging System","Energy System","PTO","PDU or HV power distribution","HV Wiring Harness","PDU Unit"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting"],"CONTACT":"Dr. Felix Richert"},"BES-C":{"SYSTEMS":["Drive System","electrical PTO","Charging System","VCU","Vehicle Control Unit","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit","Battery electric Vehicle"],"DEVDIS":["Simulation","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","HiL Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Dr. Rene Savelsberg"},"BEM":{"SYSTEMS":["mechanical PTO","E-Drive","EDU","ED Assembly","Park Lock"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting"],"CONTACT":"Thomas Flecke"},"BEM-M":{"SYSTEMS":["Axle","Gear Box or Transmission","Gears","Differential","Park Lock Actuator","E-Drive","EDU","Rotor","Stator","HV Busbars","Shafts"],"DEVDIS":["CAD Design","CAE Simulation","Stuctural Simulation","Thermal Simulation","Lifetime Simulation","Bill of Material","Patent Research","Assembly","EMC"],"CONTACT":"Dr. Gereon Hellenbroich"},"BEM-E":{"SYSTEMS":["E-Motor","Inverter or Power Electronics","DCDC power electronic","DCDC hardware","E-Motor Magnets"],"DEVDIS":["Power electronic Design","Power electronic Simulation","EMV Simulation","Electromagnetic simulation","E/E Development"],"CONTACT":"Dr. Annegret Wörndle"},"BEM-C":{"SYSTEMS":["TCU","Transmission control Unit","PLCU","Park lock control unit","DCDC software","OBC Software","Power Electronics Software","Axle"],"DEVDIS":["Simulation","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","HiL Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Dr. Jan Nowack"},"BEB":{"categories":["HV-battery system"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting"],"CONTACT":"Matthias Rudolph"},"BEB-C":{"SYSTEMS":["BMS or Battery Management System","Master BMS","Slave BMS or Slave Units"],"DEVDIS":["Simulation","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","HiL Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Lennart Bauer"},"BEB-S":{"SYSTEMS":["Battery Pack","Battery Modules","Battery Cooling System","Battery Cells","Battery Slaves","Cell Connectors"],"DEVDIS":["Simulation","Benchmarking","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Concept Development","Safety Engineering"],"CONTACT":"Dr. Andreas Averberg"},"BEB-M":{"SYSTEMS":["Battery Housing","Battery Modules","Battery Cooling System","Battery Cells","Battery Slaves","Cell Connectors","Busbars","PDU Housing"],"DEVDIS":["CAD Design","CAE Simulation","Stuctural Simulation","Thermal Simulation"],"CONTACT":"Holger Schever"},"BES-V":{"SYSTEMS":["Drive System","electrical PTO","Charging System","VCU","Vehicle Control Unit","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit","E-Motor","Park lock","Transmission","Electric Vehicle","Inverter"],"DEVDIS":["Calibration","HiL Validation","Vehicle Validation","Testing","Homologarion Support","Test Management","DVP"],"CONTACT":"Imre Pörgye"},"BEO-S":{"SYSTEMS":["Drive System","electrical PTO","Charging System","VCU","Vehicle Control Unit","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit"],"DEVDIS":["Simulation","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","HiL Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Samy Behroos"},"BMT":{"SYSTEMS":["ICE","Energy System","Electrolyzer","Test bench","Software","Measurement Equipment","Fuel Cell","Powertrain","Hybrid"],"DEVDIS":["Project management","Simulation","CAE","Testing","Consulting","Development","Design","CAD "],"CONTACT":"Dominik Lückmann"},"BMB-D":{"SYSTEMS":["ICE","Component","Energy System","Electrolyzer","Test bench","Software","Measurement Equipment","Fuel Cell","Powertrain","Hybrid"],"DEVDIS":["Project management","Simulation","CAE","Testing","Consulting","Development","Design","CAD"],"CONTACT":"Mike Souren"},"BMB-M":{"SYSTEMS":["ICE","Energy System","Electrolyzer","Fuel Cell","Powertrain","Hybrid","Turbomaschines","Large Bore Engines","Off-Road","Motorcycle","Passenger car","Heavy Duty","Medium Duty","Fuel","Gasoline","Diesel","Methanol","Ammonia","Hydrogen"],"DEVDIS":["CAE","Simulation","CAE Coordination","MBS","FEA","CFD","1D","3D","Casting Simulation","EHD","Fatigue Simulation","Vibration","Ventilation","Lubrication","Cooling","Software","Training","Workshop","CAE Assessment","Consultancy","Failure Analysis"],"CONTACT":"Sven Lauer"},"BMB-P":{"SYSTEMS":["ICE","Energy System","Electrolyzer","Test bench","Software","Measurement Equipment","Fuel Cell","Powertrain","Hybrid","Turbomaschines","Test bench design","Single Cylinder","Aerometer"],"DEVDIS":["Project management","Trouble Shooting","Measurement Equipment","Testing","Consulting","Development","Design","CAD","DVP","FMEA","Validation","Benchmarking","telemetry","products"],"CONTACT":"Max Stadermann"},"BMT-S":{"SYSTEMS":["ICE","Energy System","Fuels","Combustion","Powertrain","Hybrid","Gasoline","Diesel","Hydrogen","Methanol","Ammonia","Large Bore","Pre Ignition","Knock","EGR","Injector","Dual Fuel","Performance","Heavy-Duty","Medium-Duty","Off-Road","Motorcycle","Passenger car","Heavy Duty","Fuel","Electrolyzer","Fuel Cell","LCA","Sustainability","Heatpump","Heating System","Cooling","Function Development","HIL","Exhaust systems","EIL","SIL","MIL","Catalyst","TWC","SCR","DPF","LNT","Start","Aging","Ash Loading","GPF","Aftertreatment","RDE","Diagnostic","Diagnosis","OBM","On Board Monitoring","EVAP","Cooling","Heated Cat","Engine Cooling System","Fuel cell Cooling System","Secondary  Air","EVAP (Evaporative emissions)","Fuel System","Crankcase Ventilation","VVT (Variable Valve Timing)","EEEDWS","Drive System"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting","Simulation","Function Developemnt","Project management","Combustion Development","Thermodynamics","Troubleshooting","DVP","Systems Engineering","System FMEA","CFD","1D Simulation","3D Simulation","CAE","Lubrication","Cooling","Training","Workshop","Emission Simulation","Concept development","Requirement Management","HiL Testing","Thermal Simulation","Lifetime Simulation","System Development"],"CONTACT":"Andreas Balazs"},"BMC-O":{"SYSTEMS":["Light Duty vehicles","Medium Duty Vehicles","Heavy Duty Vehicles","ICE Gasoline conventionel","ICE Diesel conventionel","ICE Gasoline Hybrid","ICE Diesel Hybrid","Fuel Cell","Hydrogen combustion","E-Fuel","Alternative Fuels","Ethanol E0 to E100","OBD 2","OBD 1","Diagnostic","Diagnosis","OBM","On Board Monitoring","Scan tool","EVAP","Engine","Exhaust","Cooling","Comprehensive Components","Catalyst","Heated Cat","DPF Diesel Particulate Filter","GPF Gasoline Particulate Filter","Engine Cooling System","Fuel cell Cooling System","CSERS","Secondary  Air","Exhaust Gas Sensors","EVAP (Evaporative emissions)","Fuel System","CCM (incl. AC-System/DOR) comprehensive components monitoring","Other emiss. control systems","FuelcellAnode","Airpath","Crankcase Ventilation","EGR (Exhaust Gas recirculation)","Cylinder Imbalance","Misfire","VVT (Variable Valve Timing)","Fuelcell Kathode","Fuelcell Stack & System","felx fuel","SCR (Selective Catalytic Reduction)","MIL (Malfunction Indicator Lamp)","Check Engine","EEEDWS","Catalyst aging"],"DEVDIS":["Testing","Validation","Calibration","Project Management","Consulting","Concept","Legislation Assessment","Field fix","Certification","Tools","Processes","OBD Demo","Test trips","Homologation","PVE","ENDA","ECU software change accessment and migration","OBD fault management","Default Actions","Fault Reactions","Workshop Tester / End of Line","40 CFR Part 1066","Legislation interpretation","EU7","LEV IV","Tier IV","China 7"],"CONTACT":"Kortenoeven, Dennis"}}
        #                 ###

        #                 Please summarize the following RFQ text into a SUMMARY of maximum 300 words.
        #                 Then identify the categories <System 1>, <System 2>, ... that are relevant in the RFQ. Do not identify categories that are not mentioned in DEPSTRUCT and list only the most important categories. Do not list more than 5 categories.
        #                 Then identify the work activities <work activity 1>, <work activity 2>, ... in the RFQ that belong to the identified work activities (DEVDIS). Do not identify work activities that are not mentioned in DEPSTRUCT.
        #                 Afterwards identify the responsible departments <department-1>, <department-2>, ... which are responsible for both the identified categories/SYSTEMS as well as work activities (DEVDIS). Do not identify departments that are not mentioned in DEPSTRUCT.
        #                 Estimate the share of responsibility in percent of each department <responsibility-share-department-1>, <responsibility-share-department-2>, ... in the RFQ. The responsibility should represent the relevance of the RFQ for the identified departments while the sum of all responsibility shares should be equal to 100%. For example: When category "electrical PTO" is identified twice in the RFQ and "VCU" is identified twice in the RFQ and category "Battery Cooling System" once, then the responsibility share of BES-C should 75 percent and of BEB-S should be 25 percent.
        #                 Identify the contact-person of each department <contact-person 1>, <contact-person-2>. For example: "Dr. Felix Richert" is contact-person for Department "BES".
        #                 Please use valid JSON according to the following format for your response:
        #                 {"Summary": "SUMMARY","Main-Categories": {"Main Category 1": "<top-category-1>","Main Category 2": "<top-category-2>"},"Sub-Categories": {"Sub Category 1": "<work activity 1>","Sub Category 2": "<work activity 2>","Sub Category 3": "<work activity 3>"},"Departments": {"Department-1": {"Name": "<department-1>", "Share": "responsibility-share-department-1", "Contact": "contact-person-1"},"Department-2": {"Name": "<department-2>", "Share": "responsibility-share-department-2", "Contact": "contact-person-2"}},  "Project-Start": "PSTART", "Project-Duration": "PDURATION", "Proposal-Deadline": "DEADLINE", "Contact-Technical": "CONTACTTECH", "Contact-Commercial": "CONTACTCOMM"}
        #                 Example: {"Summary": "ABC Automotive GmbH is seeking a qualified engineering service provider to develop the software of an electric powertrain for a new model of electric SUV. The software development project will involve designing and implementing the software architecture and modules for the electric powertrain, integrating and testing the software with the hardware components and the vehicle platform, ensuring the software compliance with the relevant standards and regulations, and providing documentation, training, and support for the software throughout the project lifecycle. The ESPs who are interested in participating in the software development project must submit their quotations by 08/15/2024, 17:00 UTC, to the following email address: rfq@abc-automotive.com.","Main-Categories": {"Main Category 1": "VCU or Vehicle Control Unit or Powertrain Control Unit or Powertrain Domain Controller"},"Sub-Categories": {"Sub Category 1": "VCU Software","Sub Category 2": "VCU Hardware"},"Departments": {"Department-1": {"Name": "BES-C", "Share": "100 %"}}}

        #                 ### RFQ ###
        #                 """

        # Define Instructions
        instructions = """You are a Sales-Assistant of an automotive engineering service provider company trained to assign business requests (RFQs) to responsible technical departments.
                        To do so you have to:
                        1. Learn the department structure (DEPNAMES), defining what department is responsible for what categories
                        2. Learn what department is responbile for what categories (DEPSTRUCT)
                        3. Summarize the RFQ text into a SUMMARY
                        4. Identify the categories (SYSTEMS) in the RFQ text 
                        5. Identify the work activities (DEVDIS) in the RFQ text 
                        6. Identify the responsible departments 
                        7. Evaluate the share of responsibility of each department based on the occurance of the categories, the departments are responsible for
                        8. Identify the Project start if not identifiable please state that (PSTART)
                        9. Identify the Project duration if not identifiable please state that (PDURATION)
                        10. Identify the date until when the proposal shall be submitted (DEADLINE)
                        11. Identify the contact email and or person for technical questions if not identifiable please state that (CONTACTTECH)
                        12. Identify the contact email and or person for commercial questions (maybe not a Person but a generic contact mail) if not identifiable please state (CONTACTCOMM)

                        ### 1. Learn the department structure (DEPLIST) ###
                        A department name is made up of three letters, followed by a dash, follwed by another letter. example: "BES-C". 
                        "BES" is the "System Department"
                        "BEB" is the "Battery Department"
                        "BEM" is the "E-Motor and powerelectronics Department"

                        "-C": "Function Development and Controls"
                        "-C": "Systems Engineering" 
                        "-M": "Mechanical Engineering" 
                        "-V": "Validation and Verification" 
                        "-V": "Calibration" 
                        "-T": "Testing"

                        So "BES-C" is the "System Department, Function Development and Controls" and "BEM-V": is the "E-Motor and Powerelectronics Department, Calibration"
                        ###

                        ### 2. Learn what department (DEPSTRUCT) is responbile for what categories  ###
                        Learn the following responsible departments for the categories (SYSTEMS) and the work activivities (DEVDIS) and the name of the contactperson (CONTACT). the JSON-syntax is:
                        "DEPARTMENT":{"SYSTEMS":["System 1","System 1"],"DEVDIS":["work activity 1","work activity 2"],"CONTACT":"contact-person"}
                        example: "BES":{"SYSTEMS":["Powertrain and Charging System","Energy System"],"DEVDIS":["Project Management","Supplier Handling"],"CONTACT":"Dr. Felix Richert"}
                        means: "The Department BES is respondible for the categories/SYSTEMS Powertrain and Charging System and Energy System and is responible for the work activities (DEVDIS) Project Management and Supplier Handling and Dr. Felix Richert is the contact-person."
                        Here the (DEPSTRUCT):
                        {"BES":{"SYSTEMS":["Powertrain and Charging System","Energy System","PTO","PDU or HV power distribution","HV Wiring Harness","PDU Unit"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting"],"CONTACT":"Dr. Felix Richert"},"BES-C":{"SYSTEMS":["Drive System","electrical PTO","Charging System","VCU","Vehicle Control Unit (VCU)","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit","Battery electric Vehicle"],"DEVDIS":["Simulation","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","HiL Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Dr. Rene Savelsberg"},"BEM":{"SYSTEMS":["mechanical PTO","E-Drive","Electric Drive Unit (EDU)","ED Assembly","Park Lock"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting"],"CONTACT":"Thomas Flecke"},"BEM-M":{"SYSTEMS":["Axle","Gear Box or Transmission","Gears","Differential","Park Lock Actuator","E-Drive","EDU","Rotor","Stator","HV Busbars","Shafts"],"DEVDIS":["CAD Design","CAE Simulation","Stuctural Simulation","Thermal Simulation","Lifetime Simulation","Bill of Material","Patent Research","Assembly","Electro Magnetic Capability (EMC)"],"CONTACT":"Dr. Gereon Hellenbroich"},"BEM-E":{"SYSTEMS":["E-Motor","Inverter or Power Electronics","DCDC power electronic","DCDC hardware","E-Motor Magnets"],"DEVDIS":["Power electronic Design","Power electronic Simulation","Elektro Magnetische Verträglichkeit (EMV) Simulation","Electromagnetic simulation","Electric and Electronics (E/E)  Development"],"CONTACT":"Dr. Annegret Wörndle"},"BEM-C":{"SYSTEMS":["TCU","Transmission control Unit","Park lock control unit (PLCU)","DCDC software","On Board Charger (OBC) Software","Power Electronics Software","Inverter Software","Axle","Electric Drive Unit","EDU","MCU","Motor Control Unit","Actuator","PTO","Central Drive","Brake System","Steering System","E-Bike","Micromobility"],"DEVDIS":["Simulation","Benchmarking","Function Developement","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","Hardware in the Loop (HiL) Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software","A-SPICE"],"CONTACT":"Dr. Jan Nowack"},"BEB":{"categories":["HV-battery system"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting"],"CONTACT":"Matthias Rudolph"},"BEB-C":{"SYSTEMS":["BDU","Battery Disconnect Unit","Contactor","Fuse","Current Sensor","Electrical Interfaces","Cell Supervising Circuit (CSC) or Cell Monitoring Unit (CMU) or Cell Supervising Element (CSE)","Busbar"],"DEVDIS":["Design or Electrical Design","End of Line or EoL","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Commissioning","Software Development","Software Testing","HiL Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Lennart Bauer"},"BEB-S":{"SYSTEMS":["Battery Pack","Battery Modules","Battery Cooling System","Battery Cells","Battery Slaves","Cell Connectors"],"DEVDIS":["Simulation","Benchmarking","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Architecture design","Concept Development","Safety Engineering"],"CONTACT":"Dr. Andreas Averberg"},"BEB-M":{"SYSTEMS":["Battery Housing","Battery Modules","Battery Cooling System","Battery Cells","Battery Slaves","Cell Connectors","Busbars","PDU Housing"],"DEVDIS":["CAD Design","CAE Simulation","Stuctural Simulation","Thermal Simulation"],"CONTACT":"Holger Schever"},"BES-V":{"SYSTEMS":["Drive System","electrical PTO","Charging System","Vehicle Control Unit (VCU)","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit","E-Motor","Park lock","Transmission","Electric Vehicle","Inverter"],"DEVDIS":["Calibration","HiL Validation","Vehicle Validation","Testing","Homologation Support","Test Management","Design Verification Plan (DVP)"],"CONTACT":"Imre Pörgye"},"BEO-S":{"SYSTEMS":["Drive System","electrical PTO","Charging System","Vehicle Control Unit (VCU)","Powertrain Control Unit (PCU)","Powertrain Domain Controller (PDC)","VCU Software","VCU Hardware","PDU Control Unit"],"DEVDIS":["Simulation","Benchmarking","Function Developemnt","Studies","Consulting","Systems Engineering","System Development","Requirement Management","Requirements Engineering","Architecture design","Commissioning","Software Development","Software Testing","Hardware in the Loop (HIL) Testing","Software in the Loop (SIL) Testing","Model in the Loop (MIL) Testing","Software unit and component testing","Functional Safety ISO 26262","Cyber Security ISO 21434","Software unit testing","Software component testing","Concept Development","Diagnosis Software"],"CONTACT":"Samy Behroos"},"BMT":{"SYSTEMS":["ICE","Energy System","Electrolyzer","Test bench","Software","Measurement Equipment","Fuel Cell","Powertrain","Hybrid"],"DEVDIS":["Project management","Simulation","CAE","Testing","Consulting","Development","Design","CAD "],"CONTACT":"Dominik Lückmann"},"BMB-D":{"SYSTEMS":["ICE","Component","Energy System","Electrolyzer","Test bench","Software","Measurement Equipment","Fuel Cell","Powertrain","Hybrid"],"DEVDIS":["Project management","Simulation","Computer Aided Engeneering (CAE)","Testing","Consulting","Development","Design","Computer Aided Design (CAD)"],"CONTACT":"Mike Souren"},"BMB-M":{"SYSTEMS":["Internal Combution Engine (ICE)","Energy System","Electrolyzer","Fuel Cell","Powertrain","Hybrid","Turbomaschines","Large Bore Engines","Off-Road","Motorcycle","Passenger car","Heavy Duty","Medium Duty","Fuel","Gasoline","Diesel","Methanol","Ammonia","Hydrogen"],"DEVDIS":["CAE","Simulation","CAE Coordination","Multi Body Simulation (MBS)","Finite Element Analysis (FEA)","Computer Fluid Dynamics (CFD)","1D Simulation","3D Simulation","Casting Simulation","EHD","Fatigue Simulation","Vibration","Ventilation","Lubrication","Cooling","Software","Training","Workshop","CAE Assessment","Consultancy","Failure Analysis"],"CONTACT":"Sven Lauer"},"BMB-P":{"SYSTEMS":["ICE","Energy System","Electrolyzer","Test bench","Software","Measurement Equipment","Fuel Cell","Powertrain","Hybrid","Turbomaschines","Test bench design","Single Cylinder","Aerometer"],"DEVDIS":["Project management","Trouble Shooting","Measurement Equipment","Testing","Consulting","Development","Design","CAD","Design Verification Plan (DVP)","Failure Mode and Effects Analysis (FMEA)","Validation","Benchmarking","telemetry","products"],"CONTACT":"Max Stadermann"},"BMT-S":{"SYSTEMS":["ICE","Energy System","Fuels","Combustion","Powertrain","Hybrid","Gasoline","Diesel","Hydrogen","Methanol","Ammonia","Large Bore","Pre Ignition","Knock","Exhaust Gas Recirculation (EGR)","Injector","Dual Fuel","Performance","Heavy-Duty","Medium-Duty","Off-Road","Motorcycle","Passenger car","Heavy Duty","Fuel","Electrolyzer","Fuel Cell","Life Cycle Assessment (LCA)","Sustainability","Heatpump","Heating System","Cooling","Function Development","HIL","Exhaust systems","Engine In the Loop (EIL)","Software In the Loop (SIL)","Model In the Loop (MIL)","Catalyst","Three Way Catalyst (TWC)","Selective Catalytic Reduction (SCR)","Diesel Particle Filter (DPF)","Lean NOx Trap (LNT)","Start","Aging","Ash Loading","Gasoline Particulate Filter (GPF)","Aftertreatment","Real Driving Emissions (RDE)","Diagnostic","Diagnosis","OBM","On Board Monitoring","EVAP","Cooling","Heated Cat","Engine Cooling System","Fuel cell Cooling System","Secondary  Air","EVAP (Evaporative emissions)","Fuel System","Crankcase Ventilation","VVT (Variable Valve Timing)","Excess Emissions Driver Warning System (EEEDWS)","Drive System"],"DEVDIS":["Project Management","Supplier Handling","Development","Studies","Consulting","Simulation","Function Developemnt","Project management","Combustion Development","Thermodynamics","Troubleshooting","Design Verification Plan (DVP)","Systems Engineering","System Failure Mode and Effects Analysis (FMEA)","CFD","1D Simulation","3D Simulation","CAE","Lubrication","Cooling","Training","Workshop","Emission Simulation","Concept development","Requirement Management","HiL Testing","Thermal Simulation","Lifetime Simulation","System Development"],"CONTACT":"Andreas Balazs"},"BMC-O":{"SYSTEMS":["Light Duty vehicles","Medium Duty Vehicles","Heavy Duty Vehicles","ICE Gasoline conventionel","ICE Diesel conventionel","ICE Gasoline Hybrid","ICE Diesel Hybrid","Fuel Cell","Hydrogen combustion","E-Fuel","Alternative Fuels","Ethanol E0 to E100","OBD 2","OBD 1","Diagnostic","Diagnosis","OBM","On Board Monitoring","Scan tool","EVAP","Engine","Exhaust","Cooling","Comprehensive Components","Catalyst","Heated Cat","DPF Diesel Particulate Filter","GPF Gasoline Particulate Filter","Engine Cooling System","Fuel cell Cooling System","CSERS","Secondary  Air","Exhaust Gas Sensors","EVAP (Evaporative emissions)","Fuel System","CCM (incl. AC-System/DOR) comprehensive components monitoring","Other emiss. control systems","FuelcellAnode","Airpath","Crankcase Ventilation","EGR (Exhaust Gas recirculation)","Cylinder Imbalance","Misfire","VVT (Variable Valve Timing)","Fuelcell Kathode","Fuelcell Stack & System","felx fuel","SCR (Selective Catalytic Reduction)","MIL (Malfunction Indicator Lamp)","Check Engine","EEEDWS","Catalyst aging"],"DEVDIS":["Testing","Validation","Calibration","Project Management","Consulting","Concept","Legislation Assessment","Field fix","Certification","Tools","Processes","OBD Demo","Test trips","Homologation","PVE","ENDA","ECU software change accessment and migration","OBD fault management","Default Actions","Fault Reactions","Workshop Tester / End of Line","40 CFR Part 1066","Legislation interpretation","EU7","LEV IV","Tier IV","China 7"],"CONTACT":"Kortenoeven, Dennis"},"BEO-V":{"SYSTEMS":["Drive System","electrical PTO","electric Powertrain","Charging System","VCU","Vehicle Control Unit","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit","E-Motor","Park lock","Transmission","Gearbox","Hinterachsgetriebe","Getriebe Sperre","Hinterachsgetriebesperre (HAG)","Sperre","Electric Vehicle","Inverter","Fuel Management","Tank","Integrated Fuel Module (IFM)","Tankfunktionsmodul (TFM)","Industrie Kunden"],"DEVDIS":["Integration","Calibration","HiL Validation","Vehicle Validation","Testing","Absicherung","Homologation Support","Commissioning","Test Management","Design Verification Plan (DVP)","OBD","FuSi"],"CONTACT":"Benedikt Feil"},"BEO-T":{"SYSTEMS":["Drive System","electrical PTO","electric Powertrain","Charging System","On-Board-Charger (OBC)","VCU","Vehicle Control Unit","Powertrain Control Unit","Powertrain Domain Controller","VCU Software","VCU Hardware","PDU Control Unit","Battery electric Vehicle","Motor cycle"],"DEVDIS":["Integration","Validation","Verification","Testing","Absicherung","Commissioning","Test Management","Design Verification Plan (DVP)","Defect Management","Problem Management","Issue Handling"],"CONTACT":"Norman Braun"},"BVI Design":{"SYSTEMS":["Powertrain","Exhaust System","Intake System","Engine Mounts","Demonstrator vehicle","MVP","Demo vehicle","Hydrogen Tank","Tank System","Hydrogen Pipe"],"DEVDIS":["Integration","Validation","Simulation","Verification","Troubleshooting","3D Design","2D drawings","procurement","Catia","Creo","NX","Windchill","Teamcenter","gAMS"],"CONTACT":"Daniel Schrenk"},"BVI Thermal":{"SYSTEMS":["Thermal System","Air Conditioning","Cooler","Radiator","Water Pump","Heat Pump","Interior Heating","Seat heating","Surface heating","cooling circuit","A/C","Refrigerant","Cabin climatization","Heat Exchanger","HVAC","Comfort","Defrost","Demist"],"DEVDIS":["Simulation","1D Simulation","Testing","system layout","verification","benchmark","flow measurement","energy distribution","vehicle energy management"],"CONTACT":"David Hemkemeyer"},"BVI CAE":{"SYSTEMS":["engine mounts","exhaust systems","frame","chassis"],"DEVDIS":["Simulation","Durability simulation","Strength simulation","NVH simulation","benchmark"],"CONTACT":"Ralf Stienen and Jan Ophey"},"BVI Chassis":{"SYSTEMS":["Chassis","Suspension","Brake","Tire","Wheels","Damper","Spring","ABS","Kinematics"],"DEVDIS":["Testing","Validation","Tuning","benchmark","Function Test","K&C","Simulation","Multibody simulation","Ride and Handling","Vehicle Dynamics"],"CONTACT":"Christian Kuhnke and Martin Dorn"},"BVN":{"SYSTEMS":["NVH","Acoustic","Powertrain NVH","Vehicle NVH","Noise"],"DEVDIS":["Testing","Measurements"],"CONTACT":"Christoph Steffens"},"BVT":{"SYSTEMS":["Complete Vehicle","Subsystem Testing","Voice command","Charging"],"DEVDIS":["Durability run","mileage accumulation","endurance run","function tests ","comsumption tests"],"CONTACT":"Alexander Jülich and Marco Preuß"}}
                        ###

                        Please summarize the following RFQ text into a SUMMARY of maximum 300 words.
                        Then identify the categories <System 1>, <System 2>, ... that are relevant in the RFQ. Do not identify categories that are not mentioned in DEPSTRUCT and list only the most important categories. Do not list more than 5 categories.
                        Then identify the work activities <work activity 1>, <work activity 2>, ... in the RFQ that belong to the identified work activities (DEVDIS). Do not identify work activities that are not mentioned in DEPSTRUCT.
                        Afterwards identify the responsible departments <department-1>, <department-2>, ... which are responsible for both the identified categories/SYSTEMS as well as work activities (DEVDIS). Do not identify departments that are not mentioned in DEPSTRUCT.
                        Estimate the share of responsibility in percent of each department <responsibility-share-department-1>, <responsibility-share-department-2>, ... in the RFQ. The responsibility should represent the relevance of the RFQ for the identified departments while the sum of all responsibility shares should be equal to 100%. For example: When category "electrical PTO" is identified twice in the RFQ and "VCU" is identified twice in the RFQ and category "Battery Cooling System" once, then the responsibility share of BES-C should 75 percent and of BEB-S should be 25 percent.
                        Identify the contact-person of each department <contact-person 1>, <contact-person-2>. For example: "Dr. Felix Richert" is contact-person for Department "BES".
                        Please use valid JSON according to the following format for your response:
                        {"Summary": "SUMMARY","Main-Categories": {"Main Category 1": "<top-category-1>","Main Category 2": "<top-category-2>"},"Sub-Categories": {"Sub Category 1": "<work activity 1>","Sub Category 2": "<work activity 2>","Sub Category 3": "<work activity 3>"},"Departments": {"Department-1": {"Name": "<department-1>", "Share": "responsibility-share-department-1", "Contact": "contact-person-1"},"Department-2": {"Name": "<department-2>", "Share": "responsibility-share-department-2", "Contact": "contact-person-2"}},  "Project-Start": "PSTART", "Project-Duration": "PDURATION", "Proposal-Deadline": "DEADLINE", "Contact-Technical": "CONTACTTECH", "Contact-Commercial": "CONTACTCOMM"}
                        Example: {"Summary": "ABC Automotive GmbH is seeking a qualified engineering service provider to develop the software of an electric powertrain for a new model of electric SUV. The software development project will involve designing and implementing the software architecture and modules for the electric powertrain, integrating and testing the software with the hardware components and the vehicle platform, ensuring the software compliance with the relevant standards and regulations, and providing documentation, training, and support for the software throughout the project lifecycle. The ESPs who are interested in participating in the software development project must submit their quotations by 08/15/2024, 17:00 UTC, to the following email address: rfq@abc-automotive.com.","Main-Categories": {"Main Category 1": "VCU or Vehicle Control Unit or Powertrain Control Unit or Powertrain Domain Controller"},"Sub-Categories": {"Sub Category 1": "VCU Software","Sub Category 2": "VCU Hardware"},"Departments": {"Department-1": {"Name": "BES-C", "Share": "100 %"}}}

                        ### RFQ ###
                        """

        # Initialize messages array with instructions
        messages_array = [{"role": "system", "content": instructions}]

        # Get input text and add to conversation
        messages_array.append({"role": "user", "content": textRFQ})

        # Get response
        response = client.chat.completions.create(
            model=azure_oai_deployment,
            temperature=0.2,        # What sampling temperature to use, between 0 and 2. Higher values like 0.8 will make the output more random, while lower values like 0.2 will make it more focused and deterministic. We generally recommend altering this or top_p but not both.
            max_tokens=1200,        # The maximum number of tokens that can be generated in the chat completion.
            # response_format={ "type": "json_object" }, # not available for this type of model
            messages=messages_array,
            frequency_penalty=0,    # Number between -2.0 and 2.0. Positive values penalize new tokens based on their existing frequency in the text so far, decreasing the model's likelihood to repeat the same line verbatim.
            top_p=1,                # An alternative to sampling with temperature, called nucleus sampling, where the model considers the results of the tokens with top_p probability mass. So 0.1 means only the tokens comprising the top 10% probability mass are considered. We generally recommend altering this or temperature but not both.
            presence_penalty=0,     # Number between -2.0 and 2.0. Positive values penalize new tokens based on whether they appear in the text so far, increasing the model's likelihood to talk about new topics.
            seed=None,              # This feature is in Beta. If specified, our system will make a best effort to sample deterministically, such that repeated requests with the same seed and parameters should return the same result. Determinism is not guaranteed, and you should refer to the system_fingerprint response parameter to monitor changes in the backend.
            stop=None,              # Up to 4 sequences where the API will stop generating further tokens.
            stream=None,            # If set, partial message deltas will be sent, like in ChatGPT. Tokens will be sent as data-only server-sent events as they become available, with the stream terminated by a data: [DONE] message.
            tools=None,             # A list of tools the model may call. Currently, only functions are supported as a tool. Use this to provide a list of functions the model may generate JSON inputs for. A max of 128 functions are supported.
            tool_choice=None,       # Controls which (if any) function is called by the model. none means the model will not call a function and instead generates a message. auto means the model can pick between generating a message or calling a function. Specifying a particular function via {"type": "function", "function": {"name": "my_function"}} forces the model to call that function.
        )
        generated_text = response.choices[0].message.content
        completion_tokens_used = response.usage.completion_tokens
        prompt_tokens_used = response.usage.prompt_tokens
        extracted_text = generated_text

        match = re.search('```json\n(.*)```', generated_text, re.DOTALL)
        if match:
            extracted_text = match.group(1)
        
        # Return result
        return extracted_text,prompt_tokens_used,completion_tokens_used
    except Exception as ex:
        print(ex)
        return ex
